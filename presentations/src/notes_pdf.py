import sys, os, glob, textwrap
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import matplotlib
FONT=os.path.join(matplotlib.get_data_path(),"fonts/ttf/DejaVuSans.ttf")
BOLD=os.path.join(matplotlib.get_data_path(),"fonts/ttf/DejaVuSans-Bold.ttf")
pptx_path, slides_pdf, out_pdf = sys.argv[1], sys.argv[2], sys.argv[3]
work="/tmp/_nz"; os.makedirs(work, exist_ok=True)
for f in glob.glob(work+"/s-*.jpg"): os.remove(f)
os.system(f'pdftoppm -jpeg -r 120 "{slides_pdf}" {work}/s >/dev/null 2>&1')
imgs=sorted(glob.glob(work+"/s-*.jpg"))
prs=Presentation(pptx_path)
notes=[]
for sl in prs.slides:
    t=""
    if sl.has_notes_slide: t=sl.notes_slide.notes_text_frame.text.strip()
    notes.append(t)
INK=(38,57,71); TEAL=(14,140,155); AMBER=(224,138,43); GRAY=(124,139,150); LINE=(220,226,229)
CW,CH=1600,1330
pages=[]
for i,imgf in enumerate(imgs):
    note = notes[i] if i<len(notes) else ""
    canvas=Image.new("RGB",(CW,CH),"white")
    d=ImageDraw.Draw(canvas)
    sl=Image.open(imgf); sw=1360; sh=int(sl.height*sw/sl.width)
    sl=sl.resize((sw,sh)); sx=(CW-sw)//2
    canvas.paste(sl,(sx,28))
    d.rectangle([sx,28,sx+sw,28+sh],outline=LINE,width=2)
    ny=28+sh+34
    d.text((sx,ny),"TALKING POINTS",font=ImageFont.truetype(BOLD,24),fill=TEAL)
    d.text((CW-sx-120,ny),f"{i+1} / {len(imgs)}",font=ImageFont.truetype(FONT,22),fill=GRAY)
    ny+=44
    area_h=CH-ny-30
    fs=27
    while fs>=16:
        f=ImageFont.truetype(FONT,fs); fb=ImageFont.truetype(BOLD,fs)
        maxchars=int(sw/(fs*0.52))
        lines=[]
        for para in note.split("\n"):
            wrapped=textwrap.wrap(para,width=maxchars) or [""]
            lines+=wrapped
        lh=int(fs*1.34)
        if len(lines)*lh<=area_h: break
        fs-=1
    f=ImageFont.truetype(FONT,fs)
    y=ny
    for ln in lines:
        # bold the leading emphasis token before a dash/colon marker
        d.text((sx,y),ln,font=f,fill=INK); y+=int(fs*1.34)
    pages.append(canvas)
pages[0].save(out_pdf,save_all=True,append_images=pages[1:],resolution=120)
print("wrote",out_pdf,"pages",len(pages))
